import fitz
from pptx import Presentation
from pptx.util import Inches
import io

PDF = "main.pdf"
OUT = "MTP_Presentation.pptx"
ZOOM = 3.0  # ~216 dpi render for crisp slides

doc = fitz.open(PDF)
prs = Presentation()
# 4:3 slide to match Beamer page aspect (1.3333)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

mat = fitz.Matrix(ZOOM, ZOOM)
for page in doc:
    pix = page.get_pixmap(matrix=mat, alpha=False)
    img = io.BytesIO(pix.tobytes("png"))
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(img, 0, 0, width=prs.slide_width, height=prs.slide_height)

prs.save(OUT)
print("wrote", OUT, "with", len(doc), "slides")
